# server.py
from fastapi import FastAPI, BackgroundTasks, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import contextlib
import subprocess
import tempfile
import glob
import chromadb
import re
import os
from pathlib import Path
from typing import List, Optional, Any
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

import fast_fetch
import image_model

# Global ML models
models = {}

@contextlib.asynccontextmanager
async def lifespan(app: FastAPI):
    print("Loading ML models (YOLO, OCR, Embedder) into memory...")
    models["yolo"] = image_model.load_yolo(0.55)
    models["ocr"] = image_model.load_ocr()
    models["embedder"] = image_model.load_embedder()
    print("Models loaded successfully.")
    yield
    models.clear()

app = FastAPI(title="StreamStamper API", lifespan=lifespan)

# VERY IMPORTANT: This allows your React app (running on a different port) to talk to this Python server.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # In production, restrict this. Fine for hackathon.
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global status tracker so React knows what the GPU is doing
pipeline_status = {
    "status": "idle", # idle, downloading, extracting, indexing, complete, error
    "message": "Waiting for URL..."
}

class VideoRequest(BaseModel):
    url: str

class ExportResult(BaseModel):
    id: Optional[Any] = None
    text: Optional[str] = ""
    metadata: Optional[dict] = {}

class ExportRequest(BaseModel):
    results: List[ExportResult]

def extract_video_id(url: str):
    match = re.search(r"(?:v=|\/)([0-9A-Za-z_-]{11})", url)
    return match.group(1) if match else "unknown"

# --- THE BACKGROUND WORKER ---
def run_pipeline(url: str):
    global pipeline_status
    video_id = extract_video_id(url)
    
    try:
        # Step 1: Download
        pipeline_status = {"status": "downloading", "message": "Downloading video and ASR..."}
        fast_fetch.download_video(url, 16, Path("downloads"), embedder=models.get("embedder"))
        
        # Find the downloaded MP4
        downloaded_files = glob.glob("downloads/*.mp4")
        if not downloaded_files:
            raise Exception("MP4 not found after download.")
        target_mp4 = downloaded_files[0]
        
        # Step 2: Extract Frames
        pipeline_status = {"status": "extracting", "message": f"Extracting frames from {target_mp4}..."}
        subprocess.run(["python", "streamstamper.py", target_mp4, "--workers", "16"], check=True)
        
        # Step 3: Vision Indexing
        pipeline_status = {"status": "indexing", "message": "Running YOLO and PaddleOCR on GPU..."}
        image_model.process_video_visuals(video_id, models.get("yolo"), models.get("ocr"), models.get("embedder"), Path("imges"))
        
        pipeline_status = {"status": "complete", "message": "Database is ready!"}
        
    except Exception as e:
        pipeline_status = {"status": "error", "message": str(e)}


# --- API ENDPOINTS FOR REACT ---

@app.post("/api/process")
async def start_processing(req: VideoRequest):
    """React calls this when the user clicks 'Start'"""
    # Run the processing synchronously so the fetch awaits completion
    run_pipeline(req.url)
    
    if pipeline_status["status"] == "error":
        raise HTTPException(status_code=500, detail=pipeline_status["message"])
        
    return {"message": "Pipeline completed successfully"}

@app.get("/api/status")
async def get_status():
    """React polls this every 2 seconds to update the loading bar"""
    return pipeline_status

@app.get("/api/search")
async def search_db(query: str, filter_type: str = "all"):
    """React calls this when the user types in the search bar"""
    ef = models.get("embedder")
    client = chromadb.PersistentClient(path="stream_db")
    try:
        collection = client.get_collection(name="video_index", embedding_function=ef)
    except Exception:
        return {"results": []}
    
    where_clause = {}
    if filter_type == "asr":
        where_clause = {"content_type": "asr"}
    elif filter_type == "objects":
        where_clause = {"content_type": "visual_objects"}
    elif filter_type == "ocr":
        where_clause = {"content_type": "visual_ocr"}
        
    # If no filter_type provided, search both
    if where_clause:
        results = collection.query(query_texts=[query], n_results=100, where=where_clause)
    else:
        results = collection.query(query_texts=[query], n_results=100)
        
    # Format for React
    formatted_results = []
    if results and results.get('ids') and len(results['ids']) > 0:
        for i in range(len(results['ids'][0])):
            # Primary: extract from documents array
            text_val = ""
            if results.get('documents') and len(results['documents']) > 0:
                text_val = results['documents'][0][i] or ""

            # Bulletproof fallback chain: metadata["text"] -> type-specific metadata fields
            meta = results['metadatas'][0][i] if results.get('metadatas') and len(results['metadatas']) > 0 else {}
            if not text_val:
                text_val = meta.get("text", "")
            if not text_val:
                if meta.get("content_type") == "visual_ocr":
                    text_val = meta.get("ocr_text", "")
                elif meta.get("content_type") == "visual_objects":
                    text_val = meta.get("objects", "")
                    
            formatted_results.append({
                "id": results['ids'][0][i],
                "text": text_val,
                "metadata": meta
            })
        
    return {"results": formatted_results}

@app.post("/api/export/ppt")
async def export_ppt(req: ExportRequest):
    """Generate a structured study guide PowerPoint from search results."""
    prs = Presentation()

    # ── Slide 1: Title slide ──────────────────────────────────────────────────
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Automated Lecture Study Guide"
    title_slide.placeholders[1].text = "Generated by StreamStamper"

    # ── Helper ────────────────────────────────────────────────────────────────
    def fmt_ts(s):
        if s is None: return "??"
        m, sec = divmod(int(s), 60)
        return f"{m:02d}:{sec:02d}"

    # ── Group results by timestamp bucket so each slide shows ASR + OCR ───────
    # Key = (start_sec, end_sec); value = {asr: str, ocr: str}
    from collections import defaultdict
    buckets: dict = defaultdict(lambda: {"asr": "", "ocr": ""})

    for item in req.results:
        meta = item.metadata or {}
        key  = (meta.get("start_sec"), meta.get("end_sec"))
        ctype = meta.get("content_type", "")
        text  = (item.text or "").strip()
        if ctype == "asr":
            buckets[key]["asr"] = text
        elif ctype == "visual_ocr":
            buckets[key]["ocr"] = text
        elif not ctype:
            # Untyped result: put it in whichever slot is empty
            if not buckets[key]["asr"]:
                buckets[key]["asr"] = text

    # Sort by start_sec
    sorted_keys = sorted(buckets.keys(), key=lambda k: k[0] if k[0] is not None else 0)

    # ── Content slides ─────────────────────────────────────────────────────────
    content_layout = prs.slide_layouts[1]  # Title + Content
    for (start_sec, end_sec) in sorted_keys:
        data  = buckets[(start_sec, end_sec)]
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = f"{fmt_ts(start_sec)} – {fmt_ts(end_sec)}"

        tf = slide.placeholders[1].text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = f"Spoken Context (ASR): {data['asr'] or '(no spoken text)'}"
        p1.level = 1

        p2 = tf.add_paragraph()
        p2.text = f"Visual Slide Text (OCR): {data['ocr'] or '(no screen text)'}"
        p2.level = 1

    # ── Save and return ───────────────────────────────────────────────────────
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return FileResponse(
        tmp.name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="StudyGuide.pptx",
    )

@app.get("/api/timeline")
async def get_timeline():
    """React calls this to populate the 3-column chronological popup view"""
    client = chromadb.PersistentClient(path="stream_db")
    try:
        collection = client.get_collection(name="video_index")
    except Exception:
        return {"timeline": []}
    
    # Fetch everything (for hackathon limits, this is fine. For production, you'd paginate)
    all_docs = collection.get()
    
    # Sort by start_sec
    timeline = []
    for i in range(len(all_docs['ids'])):
        timeline.append({
            "text": all_docs['documents'][i],
            "metadata": all_docs['metadatas'][i]
        })
        
    timeline.sort(key=lambda x: x["metadata"]["start_sec"])
    return {"timeline": timeline}


# ─────────────────────────────────────────────────────────────────────────────
# /api/export/presentation  — Gamma-style structured study deck
# Reads ALL ASR + OCR entries from ChromaDB directly.
# Groups into 15-minute segments and produces a multi-section PPTX.
# YOLO object data is strictly excluded.
# ─────────────────────────────────────────────────────────────────────────────
SEGMENT_SECS = 15 * 60   # 15-minute chunks

def _fmt(s) -> str:
    if s is None: return "??:??"
    m, sec = divmod(int(s), 60)
    h, m   = divmod(m, 60)
    return f"{h:02d}:{m:02d}:{sec:02d}" if h else f"{m:02d}:{sec:02d}"

def _add_slide(prs, layout_idx: int, title: str, body_lines: list[str]):
    """Helper: add a slide and populate title + body text."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.word_wrap = True
    tf.text = ""                        # clear placeholder text
    from pptx.util import Pt
    for i, line in enumerate(body_lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = line
        p.space_after = Pt(6)
    return slide

@app.post("/api/export/presentation")
async def export_presentation():
    """
    Build a structured, Gamma-style educational study deck from the indexed
    lecture. Only ASR (spoken) and OCR (screen text) are included — YOLO
    object detection results are excluded entirely.
    """
    # ── 1. Fetch all ASR + OCR from ChromaDB ─────────────────────────────────
    client = chromadb.PersistentClient(path="stream_db")
    try:
        collection = client.get_collection(name="video_index")
    except Exception:
        raise HTTPException(status_code=404,
                            detail="No indexed lecture found. Run the pipeline first.")

    asr_raw = collection.get(where={"content_type": "asr"})
    ocr_raw = collection.get(where={"content_type": "visual_ocr"})

    def _rows(raw) -> list[dict]:
        rows = []
        docs  = raw.get("documents") or []
        metas = raw.get("metadatas") or []
        for doc, meta in zip(docs, metas):
            text = (doc or "").strip() or meta.get("text", "") or meta.get("ocr_text", "")
            if text:
                rows.append({"text": text, "meta": meta})
        rows.sort(key=lambda r: r["meta"].get("start_sec", 0))
        return rows

    asr_rows = _rows(asr_raw)
    ocr_rows = _rows(ocr_raw)

    # Determine video duration for labelling
    total_secs = max(
        (r["meta"].get("end_sec", 0) for r in asr_rows + ocr_rows),
        default=0,
    )
    video_id = (asr_rows[0]["meta"].get("video_id", "Unknown") if asr_rows else
                ocr_rows[0]["meta"].get("video_id", "Unknown") if ocr_rows else "Unknown")
    n_segments = max(1, int(total_secs // SEGMENT_SECS) + (1 if total_secs % SEGMENT_SECS else 0))

    def _bucket_rows(rows, seg_idx):
        lo = seg_idx * SEGMENT_SECS
        hi = lo + SEGMENT_SECS
        return [r for r in rows if lo <= r["meta"].get("start_sec", 0) < hi]

    # ── 2. Build PPTX ─────────────────────────────────────────────────────────
    prs = Presentation()
    W, H = prs.slide_width, prs.slide_height   # noqa

    # ── Slide 1: Title / Executive Summary ───────────────────────────────────
    ts = prs.slides.add_slide(prs.slide_layouts[0])
    ts.shapes.title.text = "StreamStamper Study Deck"
    ts.placeholders[1].text = (
        f"Video ID: {video_id}  |  Duration: {_fmt(total_secs)}\n\n"
        "Multi-modal synthesis of spoken lecture (ASR) and slide text (OCR).\n"
        "Object detection (YOLO) data is excluded from this study guide.\n\n"
        "Learning Objective:\n"
        "This deck provides a systematic breakdown of the professor's spoken "
        "explanations paired with the exact text written on their slides, "
        "organised into 15-minute segments for targeted revision."
    )

    # ── Slides per 15-minute segment ─────────────────────────────────────────
    for seg in range(n_segments):
        lo_sec = seg * SEGMENT_SECS
        hi_sec = lo_sec + SEGMENT_SECS
        label  = f"{_fmt(lo_sec)} – {_fmt(hi_sec)}"

        seg_asr = _bucket_rows(asr_rows, seg)
        seg_ocr = _bucket_rows(ocr_rows, seg)

        if not seg_asr and not seg_ocr:
            continue

        # ── ASR Deep-Dive slide ───────────────────────────────────────────
        if seg_asr:
            asr_lines = []
            for r in seg_asr:
                ts_label = f"[{_fmt(r['meta'].get('start_sec'))}]"
                asr_lines.append(f"{ts_label}  {r['text']}")
            _add_slide(
                prs, 1,
                f"🎙️ Spoken Content  |  {label}",
                [
                    f"Segment {seg + 1} of {n_segments}  •  ASR Transcript",
                    "",
                    *asr_lines,
                ],
            )

        # ── OCR Slide Deep-Dive ───────────────────────────────────────────
        if seg_ocr:
            ocr_lines = []
            for r in seg_ocr:
                ts_label = f"[{_fmt(r['meta'].get('start_sec'))}]"
                ocr_lines.append(f"{ts_label}  {r['text']}")
            _add_slide(
                prs, 1,
                f"📺 Slide Text  |  {label}",
                [
                    f"Segment {seg + 1} of {n_segments}  •  OCR Screen Text",
                    "",
                    *ocr_lines,
                ],
            )

        # ── Synthesis slide (only when BOTH exist) ────────────────────────
        if seg_asr and seg_ocr:
            spoken_summary  = " … ".join(r["text"][:120] for r in seg_asr[:3])
            written_summary = " … ".join(r["text"][:120] for r in seg_ocr[:3])
            _add_slide(
                prs, 1,
                f"🔗 Synthesis  |  {label}",
                [
                    "Connecting spoken context to slide content:",
                    "",
                    "🎙️ What the professor said:",
                    spoken_summary or "(no spoken content)",
                    "",
                    "📺 What was written on the slide:",
                    written_summary or "(no slide text)",
                    "",
                    "💡 Study tip: Use the OCR formula/definition above to ground "
                    "the spoken explanation into concrete, examinable terms.",
                ],
            )

    # ── Final Review slide ────────────────────────────────────────────────────
    all_ocr_text  = "  |  ".join(r["text"][:80] for r in ocr_rows[:10])
    all_asr_count = len(asr_rows)
    all_ocr_count = len(ocr_rows)
    _add_slide(
        prs, 1,
        "📋 Automated Study Guide — Full Review",
        [
            f"Total spoken segments indexed : {all_asr_count}",
            f"Total OCR slide segments      : {all_ocr_count}",
            f"Coverage                      : {_fmt(0)} – {_fmt(total_secs)}",
            "",
            "Key Slide Text (first 10 OCR entries):",
            all_ocr_text or "(no OCR data available)",
            "",
            "Revision tip: Search specific concepts in StreamStamper's "
            "ASR / OCR search to jump directly to the timestamp.",
        ],
    )

    # ── Save and return ───────────────────────────────────────────────────────
    tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".pptx")
    prs.save(tmp.name)
    tmp.close()
    return FileResponse(
        tmp.name,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="StreamStamper_StudyDeck.pptx",
    )